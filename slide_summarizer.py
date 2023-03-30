import argparse
import os

import openai
from pptx import Presentation

# If modifying these scopes, delete the file token.json.
SCOPES = ["https://www.googleapis.com/auth/presentations.readonly"]

# The ID of a sample presentation.
# PRESENTATION_ID = os.getenv("PRESENTATION_ID")


def main():
    try:
        parser = argparse.ArgumentParser(
            description="ChatGPT pptx summarizer",
            formatter_class=argparse.ArgumentDefaultsHelpFormatter,
        )
        parser.add_argument("src", help="Source pptx file")
        parser.add_argument("dest", help="Destination txt file")
        parser.add_argument(
            "-s",
            "--select",
            nargs=2,
            type=int,
            help="Select start and end slides (Inclusive)",
        )
        parser.add_argument("-k", "--key", help="Own OpenAI API key")
        args = vars(parser.parse_args())
        prs = Presentation(args["src"])
        start = 0
        end = len(prs.slides)
        if args["select"] != None:
            start = args["select"][0] - 1
            end = args["select"][1]

        slide_text = ""
        slide_list = [slide for slide in prs.slides]
        for i, slide in enumerate(slide_list[start:end]):
            slide_text += "\nSlide " + str(i + 1) + "\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
        openai.api_key = "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"

        if args["key"] != None:
            openai.api_key = args["key"]

        # System messages are used for prompting
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "Summarize this text for me"},
                {"role": "user", "content": slide_text},
            ],
        )

        print("Slide text:")
        print(slide_text)
        print("\n-----------------------------------------------------------\n")
        print("OPENAI summary:")
        print(response.choices[0].message.content)
        print("\n-----------------------------------------------------------\n")
        print("Usage: ")
        print(response.usage)

        text_file = open(args["dest"], "w")
        text_file.write(response.choices[0].message.content)
        text_file.close()

    except Exception as err:
        print(err)


if __name__ == "__main__":
    main()
